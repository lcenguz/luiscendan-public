import subprocess
import sys
import os
import pkg_resources

def instalar_dependencias():
    librerias = {
        "pymupdf": "fitz",
        "pdfplumber": "pdfplumber",
        "python-docx": "docx",
        "pandas": "pandas",
        "openpyxl": "openpyxl",
        "python-pptx": "pptx",
        "xlrd": "xlrd"
    }
    instaladas = {pkg.key for pkg in pkg_resources.working_set}
    for lib, imp in librerias.items():
        if lib not in instaladas:
            print(f"📦 Instalando {lib}...")
            subprocess.check_call([sys.executable, "-m", "pip", "install", lib])

instalar_dependencias()

# Importaciones tras instalación
import fitz
import pandas as pd
from docx import Document
from pptx import Presentation
import xml.etree.ElementTree as ET

class LectorUniversal:
    def __init__(self, ruta):
        self.ruta = ruta
        self.extension = os.path.splitext(ruta)[1].lower()

    def leer(self):
        print(f"🔍 Procesando archivo {self.extension}: {os.path.basename(self.ruta)}")
        
        try:
            if self.extension == '.pdf':
                return self._leer_pdf()
            elif self.extension in ['.xlsx', '.xls']:
                return pd.read_excel(self.ruta)
            elif self.extension == '.csv':
                return pd.read_csv(self.ruta)
            elif self.extension == '.docx':
                return self._leer_docx()
            elif self.extension == '.pptx':
                return self._leer_pptx()
            elif self.extension == '.xml':
                return self._leer_xml()
            elif self.extension == '.txt':
                with open(self.ruta, 'r', encoding='utf-8') as f:
                    return f.read()
            else:
                return f"❌ Formato {self.extension} no soportado."
        except Exception as e:
            return f"Error al procesar: {e}"

    def _leer_pdf(self):
        texto = ""
        with fitz.open(self.ruta) as doc:
            for pag in doc:
                texto += pag.get_text()
        return texto

    def _leer_docx(self):
        doc = Document(self.ruta)
        return "\n".join([para.text for para in doc.paragraphs])

    def _leer_pptx(self):
        prs = Presentation(self.ruta)
        texto = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto.append(shape.text)
        return "\n".join(texto)

    def _leer_xml(self):
        tree = ET.parse(self.ruta)
        return ET.tostring(tree.getroot(), encoding='unicode')

# --- PRUEBA ---
if __name__ == "__main__":
    # Prueba con cualquier archivo
    archivo = "datos_salud.xlsx" # Cambia esto por tu fichero
    if os.path.exists(archivo):
        lector = LectorUniversal(archivo)
        contenido = lector.leer()
        print(contenido)
    else:
        print("Asegúrate de que el archivo existe en la carpeta.")
